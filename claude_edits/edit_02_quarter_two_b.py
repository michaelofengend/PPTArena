"""
Case 7: Update Quarter Two Data (B)
- Replace all 'Q2' with 'Second Quarter' (lowercase mid-sentence)
- Slide 9: change title (shape 18) font size to 32pt
- Slide 6 (shape 26) and Slide 9 (shape 21): un-bold x-axis label text boxes
- Slide 7 (shape 15): un-bold 'YOY Activity Variance' text box
"""
import shutil
from pptx import Presentation
from pptx.util import Pt

src = '/Users/michaelofengenden/Desktop/PPTArena/Original/QuarterTwoTestB_7302020.pptx'
dst = '/Users/michaelofengenden/Desktop/PPTArena/claude_edits/QuarterTwoTestB_ClaudeEdit.pptx'

shutil.copy2(src, dst)
prs = Presentation(dst)

# Shapes where Q2 appears mid-sentence → use lowercase 'second quarter'
# All other Q2 occurrences → 'Second Quarter'
LOWERCASE_Q2 = {
    (6, 5),   # Slide 6, shape 5: '...YOY Q2 decrease...'
    (9, 4),   # Slide 9, shape 4: '...decrease in Q2 as a result...'
}

# Un-bold these shapes entirely: (slide_number, shape_id)
UNBOLD = {
    (6, 26),  # Slide 6 x-axis label
    (7, 15),  # Slide 7 'YOY Activity Variance'
    (9, 21),  # Slide 9 x-axis label
}

for slide_idx, slide in enumerate(prs.slides):
    slide_num = slide_idx + 1

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        sid = shape.shape_id
        replacement = 'second quarter' if (slide_num, sid) in LOWERCASE_Q2 else 'Second Quarter'

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                # Replace Q2
                if 'Q2' in run.text:
                    run.text = run.text.replace('Q2', replacement)

                # Un-bold
                if (slide_num, sid) in UNBOLD:
                    run.font.bold = False

        # Slide 9, shape 18: change title font size to 32pt
        if slide_num == 9 and sid == 18:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(32)
            print(f'  Slide 9, shape 18: set font size to 32pt')

prs.save(dst)
print(f'Saved: {dst}')

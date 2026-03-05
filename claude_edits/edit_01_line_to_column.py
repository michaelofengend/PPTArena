"""
Case 47: Change Line Graph to Bar Graph
Converts the line chart to a clustered column chart, preserving data and layout.
"""
import shutil
import os
from copy import deepcopy
from pptx import Presentation
from lxml import etree

C = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def cn(tag, ns=C):
    return f'{{{ns}}}{tag}'

src = '/Users/michaelofengenden/Desktop/PPTArena/Original/ChangeLineGraphtoBar_TestA.pptx'
dst = '/Users/michaelofengenden/Desktop/PPTArena/claude_edits/ChangeLineGraphtoBar_ClaudeEdit.pptx'

shutil.copy2(src, dst)
prs = Presentation(dst)
slide = prs.slides[0]

for shape in slide.shapes:
    if not shape.has_chart:
        continue

    chart_elem = shape.chart._element
    chart_node = chart_elem.find(cn('chart'))
    plot_area = chart_node.find(cn('plotArea'))
    line_chart = plot_area.find(cn('lineChart'))

    if line_chart is None:
        print('No lineChart found.')
        break

    # Build the replacement barChart element
    bar_chart = etree.Element(cn('barChart'))

    # barDir and grouping come first
    etree.SubElement(bar_chart, cn('barDir')).set('val', 'col')
    etree.SubElement(bar_chart, cn('grouping')).set('val', 'clustered')

    # Preserve varyColors
    vary = line_chart.find(cn('varyColors'))
    if vary is not None:
        bar_chart.append(deepcopy(vary))

    # Convert each series: remove marker/smooth, update spPr for fill
    for ser in line_chart.findall(cn('ser')):
        new_ser = deepcopy(ser)

        # Remove line-chart-specific elements
        for tag in ('marker', 'smooth'):
            for el in new_ser.findall(cn(tag)):
                new_ser.remove(el)

        # Replace spPr: swap line styling for solid fill
        sp_pr = new_ser.find(cn('spPr'))
        if sp_pr is not None:
            # Remove existing ln elements
            for ln in sp_pr.findall(cn('ln', A)):
                sp_pr.remove(ln)
            # Remove any existing fill
            for fill_tag in ('solidFill', 'noFill', 'gradFill', 'pattFill'):
                for el in sp_pr.findall(cn(fill_tag, A)):
                    sp_pr.remove(el)
            # Add solid fill with accent1 at position 0
            solid_fill = etree.Element(cn('solidFill', A))
            scheme_clr = etree.SubElement(solid_fill, cn('schemeClr', A))
            scheme_clr.set('val', 'accent1')
            sp_pr.insert(0, solid_fill)
            # Add a thin neutral border so bars look clean
            ln = etree.SubElement(sp_pr, cn('ln', A))
            no_fill = etree.SubElement(ln, cn('noFill', A))

        bar_chart.append(new_ser)

    # Preserve dLbls at chart level
    d_lbls = line_chart.find(cn('dLbls'))
    if d_lbls is not None:
        bar_chart.append(deepcopy(d_lbls))

    # Preserve axis references
    for ax_id in line_chart.findall(cn('axId')):
        bar_chart.append(deepcopy(ax_id))

    # Replace lineChart with barChart in the same position
    idx = list(plot_area).index(line_chart)
    plot_area.remove(line_chart)
    plot_area.insert(idx, bar_chart)

    # Fix valAx crossBetween for column charts (should be 'between')
    val_ax = plot_area.find(cn('valAx'))
    if val_ax is not None:
        cb = val_ax.find(cn('crossBetween'))
        if cb is None:
            cb = etree.SubElement(val_ax, cn('crossBetween'))
        cb.set('val', 'between')

    print('Converted lineChart → barChart (clustered column).')
    break

prs.save(dst)
print(f'Saved: {dst}')

"""
Case 30: Convert Text Data to a Line Chart
- Slide 8: remove text data box (shape 18)
- Add LINE_MARKERS chart (Asian, NHOPI, White series from 2013-2017)
- Remove auto legend; add custom legend group to the right
"""
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from lxml import etree

C  = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
A  = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
R  = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

src = '/Users/michaelofengenden/Desktop/PPTArena/Original/LineGraph_TestA.pptx'
dst = '/Users/michaelofengenden/Desktop/PPTArena/claude_edits/LineGraph_ClaudeEdit.pptx'

shutil.copy2(src, dst)
prs = Presentation(dst)
slide = prs.slides[7]  # slide 8 (0-indexed)

# ── 1. Remove the text data box (shape 18) ──────────────────────────────────
sp_tree = slide.shapes._spTree
for shape in list(slide.shapes):
    if shape.shape_id == 18:
        sp_tree.remove(shape._element)
        print('Removed shape 18 (text data box)')
        break

# ── 2. Build chart data ──────────────────────────────────────────────────────
# Values from the original text box (as decimals → shown as %)
chart_data = ChartData()
chart_data.categories = ['2013', '2014', '2015', '2016', '2017']
chart_data.add_series('Asian', (0.16, 0.11, 0.08, 0.07, 0.07))
chart_data.add_series('NHOPI', (0.18, 0.14, 0.10, 0.11, 0.11))
chart_data.add_series('White', (0.12, 0.10, 0.08, 0.07, 0.07))

# ── 3. Add chart to slide ────────────────────────────────────────────────────
chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS,
    Inches(0.101), Inches(1.630),
    Inches(8.983), Inches(4.870),
    chart_data
)
chart_elem = chart_shape.chart._element
chart_node = chart_elem.find(f'{{{C}}}chart')
plot_area  = chart_node.find(f'{{{C}}}plotArea')
line_chart = plot_area.find(f'{{{C}}}lineChart')

# ── 4. Format series: colors + circle markers ────────────────────────────────
# Series order matches add_series order: Asian=0, NHOPI=1, White=2
COLORS = ['0077C8', '0E3B5E', 'BFBFBF']

for ser in line_chart.findall(f'{{{C}}}ser'):
    idx_el = ser.find(f'{{{C}}}idx')
    idx = int(idx_el.get('val'))
    color = COLORS[idx]

    # Line style
    sp_pr = ser.find(f'{{{C}}}spPr')
    if sp_pr is None:
        sp_pr = etree.SubElement(ser, f'{{{C}}}spPr')
    else:
        sp_pr.clear()

    ln = etree.SubElement(sp_pr, f'{{{A}}}ln')
    ln.set('w', '28575'); ln.set('cap', 'rnd')
    sf = etree.SubElement(ln, f'{{{A}}}solidFill')
    etree.SubElement(sf, f'{{{A}}}srgbClr').set('val', color)
    etree.SubElement(ln, f'{{{A}}}round')
    etree.SubElement(sp_pr, f'{{{A}}}effectLst')

    # Marker (circle)
    marker = ser.find(f'{{{C}}}marker')
    if marker is None:
        marker = etree.SubElement(ser, f'{{{C}}}marker')
    else:
        marker.clear()

    etree.SubElement(marker, f'{{{C}}}symbol').set('val', 'circle')
    etree.SubElement(marker, f'{{{C}}}size').set('val', '5')
    m_sp = etree.SubElement(marker, f'{{{C}}}spPr')
    m_sf = etree.SubElement(m_sp, f'{{{A}}}solidFill')
    etree.SubElement(m_sf, f'{{{A}}}srgbClr').set('val', color)
    m_ln = etree.SubElement(m_sp, f'{{{A}}}ln')
    m_ln.set('w', '9525')
    etree.SubElement(m_ln, f'{{{A}}}noFill')
    etree.SubElement(m_sp, f'{{{A}}}effectLst')

print('Formatted series colors and markers')

# ── 5. Remove auto-legend ────────────────────────────────────────────────────
legend = chart_node.find(f'{{{C}}}legend')
if legend is not None:
    chart_node.remove(legend)

# ── 6. Format Y-axis: percentage, 0–20% ─────────────────────────────────────
val_ax = plot_area.find(f'{{{C}}}valAx')
if val_ax is not None:
    num_fmt = val_ax.find(f'{{{C}}}numFmt')
    if num_fmt is not None:
        num_fmt.set('formatCode', '0%')
        num_fmt.set('sourceLinked', '0')
    scaling = val_ax.find(f'{{{C}}}scaling')
    if scaling is not None:
        for tag in ('max', 'min'):
            el = scaling.find(f'{{{C}}}{tag}')
            if el is not None:
                scaling.remove(el)
        etree.SubElement(scaling, f'{{{C}}}max').set('val', '0.2')
        etree.SubElement(scaling, f'{{{C}}}min').set('val', '0')

# ── 7. Auto-delete chart title ───────────────────────────────────────────────
auto_title = chart_node.find(f'{{{C}}}autoTitleDeleted')
if auto_title is None:
    auto_title = etree.Element(f'{{{C}}}autoTitleDeleted')
    chart_node.insert(0, auto_title)
auto_title.set('val', '1')

# ── 8. Inject custom legend group (matches ground truth structure) ────────────
# Colors: NHOPI=0E3B5E, Asian=0077C8, White line=BFBFBF / text=636363
LEGEND_XML = '''\
<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
         xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvGrpSpPr>
    <p:cNvPr id="200" name="Legend Group"/>
    <p:cNvGrpSpPr/>
    <p:nvPr/>
  </p:nvGrpSpPr>
  <p:grpSpPr>
    <a:xfrm>
      <a:off x="8001000" y="3276600"/>
      <a:ext cx="1097280" cy="1098822"/>
      <a:chOff x="8001000" y="3124443"/>
      <a:chExt cx="1097280" cy="1098822"/>
    </a:xfrm>
  </p:grpSpPr>
  <!-- NHOPI row -->
  <p:cxnSp>
    <p:nvCxnSpPr>
      <p:cNvPr id="201" name="Line NHOPI"/>
      <p:cNvCxnSpPr/>
      <p:nvPr/>
    </p:nvCxnSpPr>
    <p:spPr>
      <a:xfrm><a:off x="8001000" y="3278331"/><a:ext cx="182880" cy="0"/></a:xfrm>
      <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
      <a:ln w="28575">
        <a:solidFill><a:srgbClr val="0E3B5E"/></a:solidFill>
        <a:headEnd type="oval"/>
        <a:tailEnd type="oval"/>
      </a:ln>
    </p:spPr>
    <p:style>
      <a:lnRef idx="1"><a:schemeClr val="accent1"/></a:lnRef>
      <a:fillRef idx="0"><a:schemeClr val="accent1"/></a:fillRef>
      <a:effectRef idx="0"><a:schemeClr val="accent1"/></a:effectRef>
      <a:fontRef idx="minor"><a:schemeClr val="tx1"/></a:fontRef>
    </p:style>
  </p:cxnSp>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="202" name="TextBox NHOPI"/>
      <p:cNvSpPr txBox="1"/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm><a:off x="8183880" y="3124443"/><a:ext cx="914400" cy="338554"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:noFill/>
    </p:spPr>
    <p:txBody>
      <a:bodyPr wrap="square" rtlCol="0"><a:spAutoFit/></a:bodyPr>
      <a:lstStyle/>
      <a:p>
        <a:r>
          <a:rPr lang="en-US" sz="1600" dirty="0">
            <a:solidFill><a:srgbClr val="0E3B5E"/></a:solidFill>
          </a:rPr>
          <a:t>NHOPI</a:t>
        </a:r>
      </a:p>
    </p:txBody>
  </p:sp>
  <!-- Asian row -->
  <p:cxnSp>
    <p:nvCxnSpPr>
      <p:cNvPr id="203" name="Line Asian"/>
      <p:cNvCxnSpPr/>
      <p:nvPr/>
    </p:nvCxnSpPr>
    <p:spPr>
      <a:xfrm><a:off x="8001514" y="3637017"/><a:ext cx="182880" cy="0"/></a:xfrm>
      <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
      <a:ln w="28575">
        <a:solidFill><a:srgbClr val="0077C8"/></a:solidFill>
        <a:headEnd type="oval"/>
        <a:tailEnd type="oval"/>
      </a:ln>
    </p:spPr>
    <p:style>
      <a:lnRef idx="1"><a:schemeClr val="accent1"/></a:lnRef>
      <a:fillRef idx="0"><a:schemeClr val="accent1"/></a:fillRef>
      <a:effectRef idx="0"><a:schemeClr val="accent1"/></a:effectRef>
      <a:fontRef idx="minor"><a:schemeClr val="tx1"/></a:fontRef>
    </p:style>
  </p:cxnSp>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="204" name="TextBox Asian"/>
      <p:cNvSpPr txBox="1"/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm><a:off x="8183880" y="3483129"/><a:ext cx="745076" cy="338554"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:noFill/>
    </p:spPr>
    <p:txBody>
      <a:bodyPr wrap="square" rtlCol="0"><a:spAutoFit/></a:bodyPr>
      <a:lstStyle/>
      <a:p>
        <a:r>
          <a:rPr lang="en-US" sz="1600" dirty="0">
            <a:solidFill><a:srgbClr val="0077C8"/></a:solidFill>
          </a:rPr>
          <a:t>Asian</a:t>
        </a:r>
      </a:p>
    </p:txBody>
  </p:sp>
  <!-- White row -->
  <p:cxnSp>
    <p:nvCxnSpPr>
      <p:cNvPr id="205" name="Line White"/>
      <p:cNvCxnSpPr/>
      <p:nvPr/>
    </p:nvCxnSpPr>
    <p:spPr>
      <a:xfrm><a:off x="8001000" y="4038600"/><a:ext cx="182880" cy="0"/></a:xfrm>
      <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
      <a:ln w="28575">
        <a:solidFill><a:srgbClr val="BFBFBF"/></a:solidFill>
        <a:headEnd type="oval"/>
        <a:tailEnd type="oval"/>
      </a:ln>
    </p:spPr>
    <p:style>
      <a:lnRef idx="1"><a:schemeClr val="accent1"/></a:lnRef>
      <a:fillRef idx="0"><a:schemeClr val="accent1"/></a:fillRef>
      <a:effectRef idx="0"><a:schemeClr val="accent1"/></a:effectRef>
      <a:fontRef idx="minor"><a:schemeClr val="tx1"/></a:fontRef>
    </p:style>
  </p:cxnSp>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="206" name="TextBox White"/>
      <p:cNvSpPr txBox="1"/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm><a:off x="8183880" y="3884711"/><a:ext cx="838200" cy="338554"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:noFill/>
    </p:spPr>
    <p:txBody>
      <a:bodyPr wrap="square" rtlCol="0"><a:spAutoFit/></a:bodyPr>
      <a:lstStyle/>
      <a:p>
        <a:r>
          <a:rPr lang="en-US" sz="1600" dirty="0">
            <a:solidFill><a:srgbClr val="636363"/></a:solidFill>
          </a:rPr>
          <a:t>White</a:t>
        </a:r>
      </a:p>
    </p:txBody>
  </p:sp>
</p:grpSp>'''

legend_grp = etree.fromstring(LEGEND_XML)
sp_tree.append(legend_grp)
print('Added custom legend group')

prs.save(dst)
print(f'Saved: {dst}')

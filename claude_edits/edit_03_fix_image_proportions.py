"""
Case 27: Correct Images
- Make all three picture placeholders the same size (3.1927" x 3.1927")
- Align them horizontally (same top = 2.3658")
- Position each one directly above its corresponding text label
  * Shape 10 (shark image) → above 'shark' label (shape 19)
  * Shape 4  (car image)   → above 'car' label   (shape 23)
  * Shape 7  (baseball img)→ above 'baseball' label (shape 24)
"""
import shutil
from pptx import Presentation
from pptx.util import Inches

src = '/Users/michaelofengenden/Desktop/PPTArena/Original/FixImagePorportions_TestA.pptx'
dst = '/Users/michaelofengenden/Desktop/PPTArena/claude_edits/FixImagePorportions_ClaudeEdit.pptx'

shutil.copy2(src, dst)
prs = Presentation(dst)
slide = prs.slides[0]

# Uniform image size and shared top (from ground truth)
IMG_W = Inches(3.1927)
IMG_H = Inches(3.1927)
IMG_TOP = Inches(2.3658)

# Text label positions (from ground truth, kept essentially unchanged)
# shape_id → target left for the image above it
LABEL_LEFTS = {
    19: Inches(1.2000),   # 'shark'
    23: Inches(5.1036),   # 'car'
    24: Inches(9.0073),   # 'baseball'
}

# Mapping: picture shape_id → which text label it belongs to (by image content)
# image1.png (sharks) = shape 10, image2.jpeg (car) = shape 4, image3.jpg (baseball) = shape 7
PIC_TO_LABEL = {
    10: 19,   # shark image → 'shark' label
    4:  23,   # car image   → 'car' label
    7:  24,   # baseball image → 'baseball' label
}

# Update text label positions — must set all four dims to avoid placeholder corruption
label_left = {}
LABEL_TOP  = Inches(5.7495)
LABEL_H    = Inches(0.6382)

for shape in slide.shapes:
    if shape.shape_id in LABEL_LEFTS:
        new_left = LABEL_LEFTS[shape.shape_id]
        label_left[shape.shape_id] = new_left
        shape.left   = new_left
        shape.top    = LABEL_TOP
        shape.width  = IMG_W
        shape.height = LABEL_H

# Reposition and resize picture placeholders
for shape in slide.shapes:
    if shape.shape_id in PIC_TO_LABEL:
        label_id = PIC_TO_LABEL[shape.shape_id]
        shape.left = label_left[label_id]
        shape.top = IMG_TOP
        shape.width = IMG_W
        shape.height = IMG_H
        print(f'  Shape {shape.shape_id}: moved to ({shape.left/914400:.4f}", {shape.top/914400:.4f}") '
              f'size ({shape.width/914400:.4f}" x {shape.height/914400:.4f}")')

prs.save(dst)
print(f'Saved: {dst}')

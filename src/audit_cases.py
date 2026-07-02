#!/usr/bin/env python3
"""
Audit the PPTArena evaluation pairs before (re-)judging.

Runs deterministic integrity checks on every case in
src/evaluation_pairs_refined.json against the actual original / ground-truth
decks, then writes a CSV of flags plus an HTML review report so a human can
click through the suspicious cases quickly. No LLM calls.

Checks per case:
  - original / ground_truth files exist and open as valid .pptx
  - ground truth actually differs from the original (a no-op GT is a broken case)
  - prompt / style_target / category present
  - "Slide N" references in the prompt or style target that exceed the deck size
  - quoted literals in the style target that never appear in the ground-truth text
    (informational — these are the rubric strings the judge scores against)
  - slide-count changes between original and GT (informational)

Usage:
    python3 src/audit_cases.py                 # checks + report, no rendering
    python3 src/audit_cases.py --render        # also convert decks to PDF for review links
    python3 src/audit_cases.py --only-flagged  # HTML report lists flagged cases only
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import html
import json
import re
import sys
from pathlib import Path

from pptx import Presentation

SCRIPT_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = SCRIPT_DIR.parent
PAIRS_PATH = SCRIPT_DIR / "evaluation_pairs_refined.json"
AUDIT_DIR = PROJECT_ROOT / "benchmark_outputs" / "audit"
PDF_DIR = SCRIPT_DIR / "work_dir" / "generated_pdfs"

SLIDE_REF_RE = re.compile(r"[Ss]lides?\s+(\d+)")
# Lookbehind avoids treating possessive apostrophes ("Team A's") as quote openers.
QUOTED_RE = re.compile(r"(?<![A-Za-z])'([^']{4,80})'")


def sha256_of(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1 << 20), b""):
            digest.update(chunk)
    return digest.hexdigest()


def deck_info(path: Path):
    """Return (slide_count, all_text_lowercase, shape_names_lowercase) or raise."""
    prs = Presentation(str(path))
    chunks: list[str] = []
    names: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name:
                names.add(re.sub(r"\s+", " ", shape.name).lower().strip())
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        chunks.append(cell.text)
    text = re.sub(r"\s+", " ", " ".join(chunks)).lower()
    return len(prs.slides), text, names


def audit_case(case: dict, render: bool) -> dict:
    name = case.get("name", "?").strip()
    flags: list[str] = []
    info: list[str] = []
    original_rel = case.get("original", "")
    gt_rel = case.get("ground_truth", "")
    original = (PROJECT_ROOT / original_rel).resolve()
    gt = (PROJECT_ROOT / gt_rel).resolve()

    if not (case.get("prompt") or "").strip():
        flags.append("empty prompt")
    if not (case.get("style_target") or "").strip():
        flags.append("empty style_target")
    if not case.get("category"):
        flags.append("missing category")

    orig_count = gt_count = None
    orig_text = gt_text = ""
    orig_names: set[str] = set()
    gt_names: set[str] = set()
    for label, path in (("original", original), ("ground_truth", gt)):
        if not path.exists():
            flags.append(f"{label} file missing")

    if original.exists():
        try:
            orig_count, orig_text, orig_names = deck_info(original)
        except Exception as exc:
            flags.append(f"original does not open: {exc}")
    if gt.exists():
        try:
            gt_count, gt_text, gt_names = deck_info(gt)
        except Exception as exc:
            flags.append(f"ground_truth does not open: {exc}")

    if original.exists() and gt.exists():
        if sha256_of(original) == sha256_of(gt):
            flags.append("ground truth is byte-identical to original (no-op case)")

    rubric_text = f"{case.get('prompt', '')} {case.get('style_target', '')}"
    refs = sorted({int(m) for m in SLIDE_REF_RE.findall(rubric_text)})
    deck_max = max(c for c in (orig_count, gt_count) if c is not None) if (orig_count or gt_count) else None
    if refs and deck_max is not None:
        over = [r for r in refs if r > deck_max]
        if over:
            flags.append(f"rubric references slide(s) {over} but decks have at most {deck_max} slides")

    missing_literals: list[str] = []
    before_refs: list[str] = []
    if gt_text:
        for literal in QUOTED_RE.findall(case.get("style_target", "")):
            normalized = re.sub(r"\s+", " ", literal).lower().strip()
            if len(normalized) < 4 or normalized in gt_text or normalized in gt_names:
                continue
            if normalized in orig_text or normalized in orig_names:
                # Present in the original but gone from the GT: the rubric is
                # quoting before-state text the edit was supposed to remove.
                before_refs.append(literal)
            else:
                missing_literals.append(literal)
    if missing_literals:
        info.append(
            "quoted rubric text not found in either deck: "
            + "; ".join(f"'{m}'" for m in missing_literals[:6])
            + (" …" if len(missing_literals) > 6 else "")
        )

    if orig_count is not None and gt_count is not None and orig_count != gt_count:
        info.append(f"slide count changes {orig_count} -> {gt_count}")

    pdf_links = {}
    if render:
        from ppt import convert_pptx_to_pdf

        for label, path in (("original", original), ("ground_truth", gt)):
            if path.exists():
                try:
                    pdf = convert_pptx_to_pdf(str(path), str(PDF_DIR))
                    if pdf:
                        pdf_links[label] = Path(pdf)
                except Exception:
                    pass

    return {
        "name": name,
        "category": ", ".join(case.get("category") or []),
        "orig_slides": orig_count,
        "gt_slides": gt_count,
        "flags": flags,
        "info": info,
        "prompt": case.get("prompt", ""),
        "style_target": case.get("style_target", ""),
        "original_rel": original_rel,
        "gt_rel": gt_rel,
        "pdf_links": pdf_links,
    }


def write_reports(results: list[dict], only_flagged: bool) -> None:
    AUDIT_DIR.mkdir(parents=True, exist_ok=True)

    csv_path = AUDIT_DIR / "audit_report.csv"
    with csv_path.open("w", newline="", encoding="utf-8") as fh:
        writer = csv.writer(fh)
        writer.writerow(["case_name", "category", "orig_slides", "gt_slides", "flags", "info"])
        for r in results:
            writer.writerow([r["name"], r["category"], r["orig_slides"], r["gt_slides"],
                             " | ".join(r["flags"]), " | ".join(r["info"])])

    listed = [r for r in results if r["flags"] or r["info"]] if only_flagged else results
    rows = []
    for r in listed:
        badge = ("<span class='flag bad'>FLAG</span>" if r["flags"]
                 else "<span class='flag warn'>REVIEW</span>" if r["info"]
                 else "<span class='flag ok'>OK</span>")
        links = []
        for label, rel in (("original", r["original_rel"]), ("ground truth", r["gt_rel"])):
            pdf = r["pdf_links"].get(label.replace(" ", "_"))
            if pdf:
                href = Path("..") / ".." / pdf.relative_to(PROJECT_ROOT)
                links.append(f"<a href='{href}' target='_blank'>{label} (pdf)</a>")
            else:
                href = Path("..") / ".." / rel
                links.append(f"<a href='{href}' target='_blank'>{label} (pptx)</a>")
        notes = "".join(f"<div class='note bad'>{html.escape(f)}</div>" for f in r["flags"])
        notes += "".join(f"<div class='note'>{html.escape(i)}</div>" for i in r["info"])
        rows.append(f"""
        <tr>
          <td class="name">{badge}<strong>{html.escape(r['name'])}</strong>
              <span class="meta">{html.escape(r['category'])} · slides {r['orig_slides']} → {r['gt_slides']}</span>
              <span class="meta">{' · '.join(links)}</span>
              {notes}</td>
          <td class="text"><b>Prompt</b> {html.escape(r['prompt'])}</td>
          <td class="text"><b>Style target</b> {html.escape(r['style_target'])}</td>
        </tr>""")

    flagged = sum(1 for r in results if r["flags"])
    review = sum(1 for r in results if r["info"] and not r["flags"])
    html_doc = f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<title>PPTArena case audit</title>
<style>
 body{{font-family:Inter,system-ui,sans-serif;margin:24px;color:#111827;font-size:14px}}
 table{{border-collapse:collapse;width:100%}}
 td{{border-bottom:1px solid #e5e7eb;padding:12px 10px;vertical-align:top}}
 .name{{width:26%}} .text{{width:37%;color:#374151;font-size:12.5px;line-height:1.5}}
 .text b{{display:block;font-size:10.5px;text-transform:uppercase;letter-spacing:.05em;color:#6b7280}}
 .meta{{display:block;color:#9ca3af;font-size:11.5px;margin-top:3px}}
 .flag{{display:inline-block;font-size:10px;font-weight:700;padding:1px 7px;border-radius:99px;margin-right:6px}}
 .flag.bad{{background:#fef2f2;color:#b91c1c;border:1px solid #fecaca}}
 .flag.warn{{background:#fffbeb;color:#b45309;border:1px solid #fde68a}}
 .flag.ok{{background:#ecfdf5;color:#047857;border:1px solid #a7f3d0}}
 .note{{margin-top:5px;font-size:12px;color:#b45309}}
 .note.bad{{color:#b91c1c;font-weight:600}}
 a{{color:#2563eb}}
</style></head><body>
<h1>PPTArena case audit</h1>
<p>{len(results)} cases · <b>{flagged} flagged</b> · {review} to review · generated by src/audit_cases.py</p>
<table>{''.join(rows)}</table>
</body></html>"""
    (AUDIT_DIR / "index.html").write_text(html_doc, encoding="utf-8")

    print(f"\n{len(results)} cases audited: {flagged} FLAG, {review} REVIEW, "
          f"{len(results) - flagged - review} OK")
    print(f"CSV:  {csv_path.relative_to(PROJECT_ROOT)}")
    print(f"HTML: {(AUDIT_DIR / 'index.html').relative_to(PROJECT_ROOT)}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Audit PPTArena evaluation pairs.")
    parser.add_argument("--render", action="store_true",
                        help="Also convert decks to PDF (LibreOffice) for click-through review links.")
    parser.add_argument("--only-flagged", action="store_true",
                        help="Only include flagged/review cases in the HTML report.")
    parser.add_argument("--cases", default=None,
                        help="Comma-separated case-name filter (substring match).")
    args = parser.parse_args()

    pairs = json.loads(PAIRS_PATH.read_text(encoding="utf-8"))
    if args.cases:
        needles = [n.strip().lower() for n in args.cases.split(",")]
        pairs = [p for p in pairs if any(n in p["name"].lower() for n in needles)]

    results = []
    for case in pairs:
        result = audit_case(case, render=args.render)
        marker = "!" if result["flags"] else ("~" if result["info"] else " ")
        print(f"[{marker}] {result['name']}")
        for f in result["flags"]:
            print(f"      FLAG: {f}")
        for i in result["info"]:
            print(f"      review: {i}")
        results.append(result)

    write_reports(results, args.only_flagged)


if __name__ == "__main__":
    main()

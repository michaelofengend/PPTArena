import pptx.oxml.simpletypes as _pptx_simpletypes
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------------------------------------------------------------------------
# Monkey patch python-pptx integer simple type parsing to tolerate float literals
# (e.g., "6096000.0" inside malformed PPTX files).
# ---------------------------------------------------------------------------
_orig_baseint_convert = _pptx_simpletypes.BaseIntType.convert_from_xml


def _baseint_convert_from_xml(cls, str_value):
    try:
        return _orig_baseint_convert(str_value)
    except ValueError:
        try:
            return int(float(str_value))
        except ValueError:
            raise


_pptx_simpletypes.BaseIntType.convert_from_xml = classmethod(_baseint_convert_from_xml)

_orig_stcoord_convert = _pptx_simpletypes.ST_Coordinate.convert_from_xml


def _stcoord_convert_from_xml(cls, str_value):
    try:
        return _orig_stcoord_convert(str_value)
    except ValueError:
        try:
            if "i" in str_value or "m" in str_value or "p" in str_value:
                return _pptx_simpletypes.ST_UniversalMeasure.convert_from_xml(str_value)
            return _pptx_simpletypes.Emu(int(float(str_value)))
        except ValueError:
            raise


_pptx_simpletypes.ST_Coordinate.convert_from_xml = classmethod(_stcoord_convert_from_xml)

def extract_text_from_shape(shape):
    """Extracts text from a shape, handling different shape types."""
    text = ""
    if shape.has_text_frame:
        text = shape.text_frame.text
    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text += cell.text_frame.text + "\t"
            text += "\n"
    return text.strip()

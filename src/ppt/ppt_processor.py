# --- ppt_processor.py ---
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pptx.oxml.simpletypes as _pptx_simpletypes
import json
import zipfile
import os
import shutil
from pathlib import Path
import base64
import re
import tempfile
import subprocess

try:
    import aspose.slides as slides
    import aspose.pydrawing as drawing
    ASPOSE_AVAILABLE = True
except ImportError:
    ASPOSE_AVAILABLE = False

# ---------------------------------------------------------------------------
# Monkey patch python-pptx integer simple type parsing to tolerate float literals
# (e.g., "6096000.0" inside malformed PPTX files).
# ---------------------------------------------------------------------------
_orig_baseint_convert = _pptx_simpletypes.BaseIntType.convert_from_xml


def _baseint_convert_from_xml(cls, str_value):
    try:
        return _orig_baseint_convert(str_value)
    except ValueError:
        try:
            return int(float(str_value))
        except ValueError:
            raise


_pptx_simpletypes.BaseIntType.convert_from_xml = classmethod(_baseint_convert_from_xml)

_orig_stcoord_convert = _pptx_simpletypes.ST_Coordinate.convert_from_xml


def _stcoord_convert_from_xml(cls, str_value):
    try:
        return _orig_stcoord_convert(str_value)
    except ValueError:
        try:
            if "i" in str_value or "m" in str_value or "p" in str_value:
                return _pptx_simpletypes.ST_UniversalMeasure.convert_from_xml(str_value)
            return _pptx_simpletypes.Emu(int(float(str_value)))
        except ValueError:
            raise


_pptx_simpletypes.ST_Coordinate.convert_from_xml = classmethod(_stcoord_convert_from_xml)

def extract_text_from_shape(shape):
    """Extracts text from a shape, handling different shape types."""
    text = ""
    if shape.has_text_frame:
        text = shape.text_frame.text
    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text += cell.text_frame.text + "\t"
            text += "\n"
    return text.strip()

def pptx_to_json(filepath):
    """
    Converts a .pptx file to a comprehensive JSON representation.
    Captures all edit types: Content, Layout, Styling, Interactivity, and Structure.
    """
    try:
        prs = Presentation(filepath)
        
        # Presentation-level metadata
        presentation_data = {
            "filename": os.path.basename(filepath),
            "slide_width": prs.slide_width.pt if hasattr(prs, 'slide_width') else None,
            "slide_height": prs.slide_height.pt if hasattr(prs, 'slide_height') else None,
            "slides": []
        }
        
        # LOW PRIORITY: Custom properties/metadata (existence check)
        try:
            if hasattr(prs, 'core_properties'):
                props = prs.core_properties
                custom_props = {}
                if hasattr(props, 'author') and props.author:
                    custom_props["author"] = props.author
                if hasattr(props, 'title') and props.title:
                    custom_props["title"] = props.title
                if hasattr(props, 'subject') and props.subject:
                    custom_props["subject"] = props.subject
                if hasattr(props, 'keywords') and props.keywords:
                    custom_props["keywords"] = props.keywords
                if custom_props:
                    presentation_data["custom_properties"] = custom_props
        except:
            pass
        
        for i, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": i + 1,
                "shapes": [],
                "notes": "",
                # Layout & Structure info
                "slide_layout": slide.slide_layout.name if hasattr(slide, 'slide_layout') else None,
                "slide_id": slide.slide_id if hasattr(slide, 'slide_id') else None,
            }
            
            # Background info (Styling category)
            try:
                if hasattr(slide, 'background'):
                    bg = slide.background
                    slide_data["background"] = {
                        "fill_type": str(bg.fill.type) if hasattr(bg, 'fill') else None,
                    }
                    try:
                        if hasattr(bg.fill, 'fore_color') and bg.fill.fore_color.rgb:
                            slide_data["background"]["color_rgb"] = str(bg.fill.fore_color.rgb)
                    except:
                        pass
            except Exception:
                pass
            
            for shape_idx, shape in enumerate(slide.shapes):
                shape_info = {
                    "shape_id": shape.shape_id if hasattr(shape, 'shape_id') else None,
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "text": extract_text_from_shape(shape),
                    # Layout info (position, size, rotation)
                    "left": shape.left.pt if hasattr(shape, 'left') and shape.left is not None else None,
                    "top": shape.top.pt if hasattr(shape, 'top') and shape.top is not None else None,
                    "width": shape.width.pt if hasattr(shape, 'width') and shape.width is not None else None,
                    "height": shape.height.pt if hasattr(shape, 'height') and shape.height is not None else None,
                    "rotation": shape.rotation if hasattr(shape, 'rotation') else None,
                    "z_order": shape_idx,  # Z-order based on iteration order
                }
                
                # Placeholder info (Layout category)
                try:
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        shape_info["is_placeholder"] = True
                        shape_info["placeholder_type"] = str(shape.placeholder_format.type) if hasattr(shape, 'placeholder_format') else None
                except:
                    pass
                
                # Shape styling (fills, lines, effects)
                try:
                    if hasattr(shape, 'fill'):
                        fill_info = {
                            "type": str(shape.fill.type) if shape.fill else None,
                        }
                        try:
                            if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color and shape.fill.fore_color.rgb:
                                fill_info["color_rgb"] = str(shape.fill.fore_color.rgb)
                        except:
                            pass
                        
                        # HIGH PRIORITY: Gradient fills
                        try:
                            if hasattr(shape.fill, 'gradient_angle'):
                                fill_info["gradient_angle"] = shape.fill.gradient_angle
                            if hasattr(shape.fill, 'gradient_stops'):
                                stops = []
                                for stop in shape.fill.gradient_stops:
                                    stop_info = {
                                        "position": stop.position if hasattr(stop, 'position') else None,
                                    }
                                    try:
                                        if hasattr(stop, 'color') and stop.color and stop.color.rgb:
                                            stop_info["color_rgb"] = str(stop.color.rgb)
                                    except:
                                        pass
                                    stops.append(stop_info)
                                if stops:
                                    fill_info["gradient_stops"] = stops
                        except:
                            pass
                        
                        shape_info["fill"] = fill_info
                except:
                    pass
                
                try:
                    if hasattr(shape, 'line'):
                        line_info = {
                            "width": shape.line.width.pt if shape.line.width else None,
                        }
                        try:
                            if hasattr(shape.line, 'color') and shape.line.color and shape.line.color.rgb:
                                line_info["color_rgb"] = str(shape.line.color.rgb)
                        except:
                            pass
                        # HIGH PRIORITY: Line dash style
                        try:
                            if hasattr(shape.line, 'dash_style') and shape.line.dash_style:
                                line_info["dash_style"] = str(shape.line.dash_style)
                        except:
                            pass
                        shape_info["line"] = line_info
                except:
                    pass
                
                # Shadow effect
                try:
                    if hasattr(shape, 'shadow') and shape.shadow and shape.shadow.inherit:
                        shape_info["has_shadow"] = True
                except:
                    pass
                
                # Hyperlinks (Interactivity category)
                try:
                    if hasattr(shape, 'click_action') and shape.click_action and shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                        shape_info["hyperlink"] = shape.click_action.hyperlink.address
                except:
                    pass
                
                # Extract detailed TYPOGRAPHY from text frames (Content category)
                if shape.has_text_frame:
                    shape_info["paragraphs"] = []
                    for para in shape.text_frame.paragraphs:
                        para_info = {
                            "text": para.text,
                            "level": para.level,
                            "alignment": str(para.alignment) if hasattr(para, 'alignment') else None,
                            "line_spacing": para.line_spacing if hasattr(para, 'line_spacing') else None,
                            "space_before": para.space_before.pt if hasattr(para, 'space_before') and para.space_before else None,
                            "space_after": para.space_after.pt if hasattr(para, 'space_after') and para.space_after else None,
                            "runs": []
                        }
                        
                        # HIGH PRIORITY: Bullet/Numbering formats
                        try:
                            if hasattr(para, 'bullet_format') and para.bullet_format:
                                bullet_info = {}
                                if hasattr(para.bullet_format, 'type') and para.bullet_format.type:
                                    bullet_info["type"] = str(para.bullet_format.type)
                                if hasattr(para.bullet_format, 'char') and para.bullet_format.char:
                                    bullet_info["char"] = para.bullet_format.char
                                if hasattr(para.bullet_format, 'start_value') and para.bullet_format.start_value:
                                    bullet_info["start_value"] = para.bullet_format.start_value
                                if bullet_info:
                                    para_info["bullet_format"] = bullet_info
                        except:
                            pass
                        
                        for run in para.runs:
                            run_info = {
                                "text": run.text,
                                "font_name": run.font.name if run.font.name else None,
                                "font_size": run.font.size.pt if run.font.size else None,
                                "bold": run.font.bold,
                                "italic": run.font.italic,
                                "underline": run.font.underline,
                            }
                            
                            # Color info
                            try:
                                if run.font.color and run.font.color.rgb:
                                    run_info["color_rgb"] = str(run.font.color.rgb)
                            except Exception:
                                pass
                            
                            # Hyperlink in run
                            try:
                                if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                                    run_info["hyperlink"] = run.hyperlink.address
                            except:
                                pass
                            
                            para_info["runs"].append(run_info)
                        
                        shape_info["paragraphs"].append(para_info)
                
                # TABLES (Content category)
                if shape.has_table:
                    table_info = {
                        "rows": len(shape.table.rows),
                        "columns": len(shape.table.columns),
                        "cells": []
                    }
                    
                    for row_idx, row in enumerate(shape.table.rows):
                        table_info["row_height_" + str(row_idx)] = row.height.pt if row.height else None
                        for col_idx, cell in enumerate(row.cells):
                            cell_info = {
                                "row": row_idx,
                                "col": col_idx,
                                "text": cell.text_frame.text if hasattr(cell, 'text_frame') else "",
                                "paragraphs": []
                            }
                            
                            # Cell fill
                            try:
                                if hasattr(cell, 'fill') and cell.fill:
                                    cell_info["fill_type"] = str(cell.fill.type)
                                    if hasattr(cell.fill, 'fore_color') and cell.fill.fore_color and cell.fill.fore_color.rgb:
                                        cell_info["fill_color_rgb"] = str(cell.fill.fore_color.rgb)
                            except:
                                pass
                            
                            # MEDIUM PRIORITY: Table cell borders
                            try:
                                if hasattr(cell, 'borders'):
                                    borders_info = {}
                                    for side in ['top', 'bottom', 'left', 'right']:
                                        if hasattr(cell.borders, side):
                                            border = getattr(cell.borders, side)
                                            if border:
                                                border_info = {}
                                                try:
                                                    if hasattr(border, 'width') and border.width:
                                                        border_info["width"] = border.width.pt
                                                except:
                                                    pass
                                                try:
                                                    if hasattr(border, 'color') and border.color and border.color.rgb:
                                                        border_info["color_rgb"] = str(border.color.rgb)
                                                except:
                                                    pass
                                                try:
                                                    if hasattr(border, 'dash_style'):
                                                        border_info["dash_style"] = str(border.dash_style)
                                                except:
                                                    pass
                                                if border_info:
                                                    borders_info[side] = border_info
                                    if borders_info:
                                        cell_info["borders"] = borders_info
                            except:
                                pass
                            
                            # Cell typography
                            if hasattr(cell, 'text_frame'):
                                for para in cell.text_frame.paragraphs:
                                    para_info = {
                                        "text": para.text,
                                        "runs": []
                                    }
                                    
                                    for run in para.runs:
                                        run_info = {
                                            "text": run.text,
                                            "font_name": run.font.name if run.font.name else None,
                                            "font_size": run.font.size.pt if run.font.size else None,
                                            "bold": run.font.bold,
                                            "italic": run.font.italic,
                                        }
                                        try:
                                            if run.font.color and run.font.color.rgb:
                                                run_info["color_rgb"] = str(run.font.color.rgb)
                                        except:
                                            pass
                                        para_info["runs"].append(run_info)
                                    
                                    cell_info["paragraphs"].append(para_info)
                            
                            table_info["cells"].append(cell_info)
                    
                    # Column widths
                    for col_idx, col in enumerate(shape.table.columns):
                        table_info["column_width_" + str(col_idx)] = col.width.pt if col.width else None
                    
                    shape_info["table"] = table_info
                
                # IMAGES & PICTURES (Content category)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    shape_info["is_picture"] = True
                    try:
                        if hasattr(shape, 'image'):
                            shape_info["image_type"] = shape.image.content_type if hasattr(shape.image, 'content_type') else None
                            shape_info["image_filename"] = shape.image.filename if hasattr(shape.image, 'filename') else None
                    except:
                        pass
                
                # CHARTS (Content category)
                if hasattr(shape, 'has_chart') and shape.has_chart:
                    shape_info["is_chart"] = True
                    try:
                        chart = shape.chart
                        chart_info = {
                            "chart_type": str(chart.chart_type) if hasattr(chart, 'chart_type') else None,
                            "has_title": chart.has_title if hasattr(chart, 'has_title') else None,
                            "chart_title": chart.chart_title.text_frame.text if hasattr(chart, 'chart_title') and chart.has_title else None,
                        }
                        
                        # HIGH PRIORITY: Chart series data
                        try:
                            if hasattr(chart, 'series'):
                                series_list = []
                                for series in chart.series:
                                    series_info = {
                                        "name": series.name if hasattr(series, 'name') else None,
                                    }
                                    # Extract data values if available
                                    try:
                                        if hasattr(series, 'values') and series.values:
                                            series_info["values"] = list(series.values)
                                    except:
                                        pass
                                    series_list.append(series_info)
                                if series_list:
                                    chart_info["series"] = series_list
                        except:
                            pass
                        
                        # Extract categories if available
                        try:
                            if hasattr(chart, 'categories') and chart.categories:
                                chart_info["categories"] = list(chart.categories)
                        except:
                            pass
                        
                        shape_info["chart"] = chart_info
                    except:
                        pass
                
                # GROUP SHAPES (Layout category)
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    shape_info["is_group"] = True
                    shape_info["group_shape_count"] = len(shape.shapes) if hasattr(shape, 'shapes') else 0
                
                # HIGH PRIORITY: Shape connectors (for flowcharts, org charts)
                try:
                    if hasattr(shape, 'connector_format'):
                        connector_info = {}
                        if hasattr(shape.connector_format, 'begin_connected') and shape.connector_format.begin_connected:
                            connector_info["begin_connected"] = True
                            if hasattr(shape.connector_format, 'begin_connection_site'):
                                connector_info["begin_connection_site"] = shape.connector_format.begin_connection_site
                        if hasattr(shape.connector_format, 'end_connected') and shape.connector_format.end_connected:
                            connector_info["end_connected"] = True
                            if hasattr(shape.connector_format, 'end_connection_site'):
                                connector_info["end_connection_site"] = shape.connector_format.end_connection_site
                        if connector_info:
                            shape_info["connector"] = connector_info
                except:
                    pass
                
                # Alt Text / Accessibility (Structure category)
                try:
                    if hasattr(shape, 'name'):
                        shape_info["alt_text"] = shape.name
                except:
                    pass
                
                slide_data["shapes"].append(shape_info)
            
            # NOTES (Structure category)
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                slide_data["notes"] = text_frame.text.strip()
            
            # SLIDE TRANSITIONS (Interactivity category) - LOW PRIORITY: Enhanced detection
            try:
                # Check for transition via element property
                if hasattr(slide, 'element') and hasattr(slide.element, 'transition'):
                    transition = slide.element.transition
                    if transition is not None:
                        slide_data["has_transition"] = True
                        # Try to extract transition type if available
                        try:
                            if hasattr(transition, 'type'):
                                slide_data["transition_type"] = str(transition.type)
                        except:
                            pass
                elif hasattr(slide, 'slide_transitions'):
                    slide_data["has_transition"] = True
            except:
                pass
            
            # LOW PRIORITY: SmartArt detection (existence check)
            try:
                has_smartart = False
                for shape in slide.shapes:
                    # SmartArt is typically detected as a group with diagram data
                    if hasattr(shape, 'element') and hasattr(shape.element, 'tag'):
                        tag = str(shape.element.tag)
                        if 'dgm' in tag.lower() or 'diagram' in tag.lower():
                            has_smartart = True
                            break
                if has_smartart:
                    slide_data["has_smartart"] = True
            except:
                pass
            
            presentation_data["slides"].append(slide_data)
        
        return presentation_data
    except Exception as e:
        print(f"Error converting {filepath} to JSON: {e}")
        raise

def diff_pptx_json(ground_truth_json, prediction_json, initial_json=None):
    """
    Performs a deep comparison between ground_truth and prediction JSON structures.
    Returns a structured diff with only the differences, organized by slide and shape.
    
    Returns:
        {
            "has_differences": bool,
            "similarity_score": float (0-1),
            "differences": [
                {
                    "slide_number": int,
                    "shape_id": str,
                    "shape_name": str,
                    "property_path": str,
                    "ground_truth_value": any,
                    "prediction_value": any,
                    "initial_value": any (if provided),
                    "change_type": "added" | "removed" | "modified"
                }
            ]
        }
    """
    differences = []
    
    def normalize_value(val):
        """Normalize values for comparison (handle floats, None, etc.)"""
        if val is None:
            return None
        if isinstance(val, float):
            return round(val, 2)  # Round to 2 decimals for comparison
        if isinstance(val, (list, tuple)):
            return [normalize_value(v) for v in val]
        if isinstance(val, dict):
            return {k: normalize_value(v) for k, v in val.items()}
        return val
    
    def values_match(val1, val2, tolerance=0.01):
        """Check if two values match, with tolerance for floats (1% tolerance)"""
        val1 = normalize_value(val1)
        val2 = normalize_value(val2)
        
        if val1 == val2:
            return True
        
        # Handle float comparisons with tolerance
        if isinstance(val1, (int, float)) and isinstance(val2, (int, float)):
            if val1 == 0 and val2 == 0:
                return True
            if val1 == 0 or val2 == 0:
                return abs(val1 - val2) < 0.01
            return abs(val1 - val2) / max(abs(val1), abs(val2)) < tolerance
        
        return False
    
    def compare_dict(path, gt_dict, pred_dict, init_dict=None):
        """Recursively compare two dictionaries"""
        all_keys = set(gt_dict.keys()) | set(pred_dict.keys())
        
        for key in all_keys:
            current_path = f"{path}.{key}" if path else key
            
            gt_val = gt_dict.get(key)
            pred_val = pred_dict.get(key)
            init_val = init_dict.get(key) if init_dict else None
            
            # Skip certain keys that are not semantic (auto-generated IDs)
            # Note: z_order is now included as it's semantically important for layering
            if key in ['shape_id', 'slide_id']:
                continue
            
            if key not in pred_dict:
                differences.append({
                    "property_path": current_path,
                    "change_type": "removed",
                    "ground_truth_value": gt_val,
                    "prediction_value": None,
                    "initial_value": init_val
                })
            elif key not in gt_dict:
                differences.append({
                    "property_path": current_path,
                    "change_type": "added",
                    "ground_truth_value": None,
                    "prediction_value": pred_val,
                    "initial_value": init_val
                })
            elif isinstance(gt_val, dict) and isinstance(pred_val, dict):
                compare_dict(current_path, gt_val, pred_val, init_val if isinstance(init_val, dict) else None)
            elif isinstance(gt_val, list) and isinstance(pred_val, list):
                compare_list(current_path, gt_val, pred_val, init_val if isinstance(init_val, list) else None)
            elif not values_match(gt_val, pred_val):
                differences.append({
                    "property_path": current_path,
                    "change_type": "modified",
                    "ground_truth_value": gt_val,
                    "prediction_value": pred_val,
                    "initial_value": init_val
                })
    
    def compare_list(path, gt_list, pred_list, init_list=None):
        """Compare two lists"""
        # For lists of primitives, do simple comparison
        if gt_list and not isinstance(gt_list[0], (dict, list)):
            if not values_match(gt_list, pred_list):
                differences.append({
                    "property_path": path,
                    "change_type": "modified",
                    "ground_truth_value": gt_list,
                    "prediction_value": pred_list,
                    "initial_value": init_list
                })
        else:
            # For lists of dicts (like paragraphs, runs, shapes), compare element by element
            max_len = max(len(gt_list), len(pred_list))
            for i in range(max_len):
                gt_item = gt_list[i] if i < len(gt_list) else None
                pred_item = pred_list[i] if i < len(pred_list) else None
                init_item = init_list[i] if init_list and i < len(init_list) else None
                
                current_path = f"{path}[{i}]"
                
                if gt_item is None:
                    differences.append({
                        "property_path": current_path,
                        "change_type": "added",
                        "ground_truth_value": None,
                        "prediction_value": pred_item,
                        "initial_value": init_item
                    })
                elif pred_item is None:
                    differences.append({
                        "property_path": current_path,
                        "change_type": "removed",
                        "ground_truth_value": gt_item,
                        "prediction_value": None,
                        "initial_value": init_item
                    })
                elif isinstance(gt_item, dict) and isinstance(pred_item, dict):
                    compare_dict(current_path, gt_item, pred_item, init_item if isinstance(init_item, dict) else None)
    
    # Compare slides
    gt_slides = ground_truth_json.get("slides", [])
    pred_slides = prediction_json.get("slides", [])
    init_slides = initial_json.get("slides", []) if initial_json else []
    
    for slide_idx in range(max(len(gt_slides), len(pred_slides))):
        gt_slide = gt_slides[slide_idx] if slide_idx < len(gt_slides) else {}
        pred_slide = pred_slides[slide_idx] if slide_idx < len(pred_slides) else {}
        init_slide = init_slides[slide_idx] if slide_idx < len(init_slides) else {}
        
        slide_number = slide_idx + 1
        
        # Compare slide-level properties
        slide_diff_start = len(differences)
        compare_dict(f"slide_{slide_number}", gt_slide, pred_slide, init_slide)
        
        # Tag differences with slide context
        for diff in differences[slide_diff_start:]:
            diff["slide_number"] = slide_number
            
            # Extract shape context if available
            if "shapes[" in diff["property_path"]:
                shape_idx_match = re.search(r'shapes\[(\d+)\]', diff["property_path"])
                if shape_idx_match:
                    shape_idx = int(shape_idx_match.group(1))
                    shapes = pred_slide.get("shapes", [])
                    if shape_idx < len(shapes):
                        diff["shape_name"] = shapes[shape_idx].get("name", f"Shape_{shape_idx}")
                        diff["shape_type"] = shapes[shape_idx].get("type", "Unknown")
    
    # Calculate similarity score
    total_properties = len(differences) + 100  # Baseline to avoid division by zero
    similarity_score = 1.0 - (len(differences) / total_properties)
    
    return {
        "has_differences": len(differences) > 0,
        "similarity_score": max(0.0, min(1.0, similarity_score)),
        "total_differences": len(differences),
        "differences": differences
    }


def should_use_xml_fallback(instruction_text, has_smartart=False, has_transitions=False):
    """
    Determines if XML-based analysis should be used instead of JSON.
    Returns True if the instruction involves SmartArt, animations, or complex transitions.
    """
    if not instruction_text:
        return False
    
    instruction_lower = instruction_text.lower()
    
    # Keywords that indicate need for XML analysis
    xml_keywords = [
        'smartart', 'smart art', 'diagram',
        'animation', 'animate', 'motion path', 'entrance', 'exit', 'emphasis',
        'transition', 'slide transition', 'fade', 'wipe', 'dissolve'
    ]
    
    for keyword in xml_keywords:
        if keyword in instruction_lower:
            return True
    
    # Also check if slides actually have SmartArt or transitions
    if has_smartart or has_transitions:
        return True
    
    return False


def format_diff_for_judge(diff_result, instruction_text=""):
    """
    Formats the diff result into a human-readable format optimized for LLM judge.
    Groups changes by slide and shape, provides context, and highlights key differences.
    """
    if not diff_result.get("has_differences"):
        return "✓ No differences detected between prediction and ground truth. Both presentations are identical."
    
    differences = diff_result.get("differences", [])
    if not differences:
        return "✓ No significant differences detected."
    
    # Group differences by slide and shape
    by_slide = {}
    for diff in differences:
        slide_num = diff.get("slide_number", 0)
        if slide_num not in by_slide:
            by_slide[slide_num] = {}
        
        shape_name = diff.get("shape_name", "Slide-level property")
        if shape_name not in by_slide[slide_num]:
            by_slide[slide_num][shape_name] = []
        
        by_slide[slide_num][shape_name].append(diff)
    
    # Format output
    lines = [
        f"📊 Comparison Summary: {diff_result['total_differences']} differences found",
        f"Similarity Score: {diff_result['similarity_score']:.1%}",
        ""
    ]
    
    # Add instruction context if provided
    if instruction_text:
        lines.append(f"Instruction: {instruction_text}")
        lines.append("")
    
    lines.append("Detailed Differences:")
    lines.append("=" * 80)
    
    for slide_num in sorted(by_slide.keys()):
        lines.append(f"\n📄 Slide {slide_num}:")
        
        for shape_name, shape_diffs in by_slide[slide_num].items():
            lines.append(f"  └─ {shape_name}")
            
            for diff in shape_diffs:
                prop_path = diff["property_path"]
                # Simplify path for readability
                simple_path = prop_path.split('.')[-1] if '.' in prop_path else prop_path
                simple_path = re.sub(r'\[\d+\]', '', simple_path)  # Remove array indices
                
                change_type = diff["change_type"]
                gt_val = diff.get("ground_truth_value")
                pred_val = diff.get("prediction_value")
                init_val = diff.get("initial_value")
                
                # Format based on change type
                if change_type == "modified":
                    status = "✗" if gt_val != pred_val else "✓"
                    lines.append(f"      {status} {simple_path}:")
                    lines.append(f"         Ground Truth: {gt_val}")
                    lines.append(f"         Prediction:   {pred_val}")
                    if init_val is not None and init_val != gt_val:
                        lines.append(f"         Initial:      {init_val}")
                elif change_type == "added":
                    lines.append(f"      + {simple_path}: {pred_val} (added in prediction)")
                elif change_type == "removed":
                    lines.append(f"      - {simple_path}: {gt_val} (missing in prediction)")
    
    return "\n".join(lines)


def extract_xml_from_pptx(pptx_filepath, output_folder):
    """
    Extracts all constituent XML files from a .pptx file.
    Returns a list of full paths to the extracted XML files.
    """
    extracted_files_paths = []
    try:
        Path(output_folder).mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(pptx_filepath, 'r') as pptx_zip:
            for member_info in pptx_zip.infolist():
                member_name = member_info.filename
                if not member_info.is_dir() and member_name.endswith(('.xml', '.rels')):
                    target_path = os.path.join(output_folder, member_name)
                    os.makedirs(os.path.dirname(target_path), exist_ok=True)
                    with pptx_zip.open(member_name) as source, open(target_path, "wb") as target:
                        shutil.copyfileobj(source, target)
                    extracted_files_paths.append(target_path)
        return extracted_files_paths
    except Exception as e:
        print(f"Error extracting XML from {pptx_filepath}: {e}")
        raise

def create_modified_pptx(original_pptx_path, modified_xml_map, output_pptx_path):
    """
    Creates a new .pptx file by taking an original .pptx, and either updating
    existing internal XML files or adding new ones based on the modified_xml_map.
    """
    temp_output_pptx_path = output_pptx_path + ".tmp"
    try:
        os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
        
        # Keep track of which files from the modification map we have used.
        processed_map_files = set()

        with zipfile.ZipFile(original_pptx_path, 'r') as zin:
            with zipfile.ZipFile(temp_output_pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
                # 1. Iterate through existing files in the original PPTX.
                for item in zin.infolist():
                    item_name_normalized = item.filename.replace("\\", "/")
                    if item_name_normalized in modified_xml_map:
                        # If an existing file is in our map, write the modified content.
                        new_content = modified_xml_map[item_name_normalized]
                        zout.writestr(item, new_content.encode('utf-8'))
                        processed_map_files.add(item_name_normalized)
                    else:
                        # Otherwise, copy the original file as-is.
                        buffer = zin.read(item.filename)
                        zout.writestr(item, buffer)
                
                # 2. Iterate through the map to find any *new* files to add.
                for filename, content in modified_xml_map.items():
                    if filename not in processed_map_files:
                        print(f"Adding new file to PPTX archive: {filename}")
                        # writestr can take a ZipInfo object or a string filename.
                        zout.writestr(filename, content.encode('utf-8'))

        os.replace(temp_output_pptx_path, output_pptx_path)
        print(f"Modified PPTX successfully created at: {output_pptx_path}")
        return True
    except Exception as e:
        print(f"Error creating modified PPTX at {output_pptx_path}: {e}")
        if os.path.exists(temp_output_pptx_path):
            os.remove(temp_output_pptx_path)
        return False

def export_slides_to_images(pptx_filepath, output_folder):
    """
    Converts each slide of a .pptx file to PNG images using LibreOffice if
    available. Returns a list of file paths to the generated images. If
    conversion fails, an empty list is returned.
    """
    # Try to find a valid executable for LibreOffice
    libreoffice_exec = None
    for name in ['soffice', 'libreoffice']:
        if shutil.which(name):
            libreoffice_exec = name
            break

    if not libreoffice_exec:
        print("ERROR: Could not find 'soffice' or 'libreoffice' command.")
        print("Please ensure LibreOffice is installed and that its program directory is in your system's PATH.")
        return []

    # Use abspath to handle relative paths correctly
    abs_pptx_filepath = os.path.abspath(pptx_filepath)
    abs_output_folder = os.path.abspath(output_folder)
    Path(abs_output_folder).mkdir(parents=True, exist_ok=True)

    try:
        # On macOS, the command might be different, but 'libreoffice' is standard for PATH
        # The user might need to symlink /Applications/LibreOffice.app/Contents/MacOS/soffice
        cmd = [
            libreoffice_exec,
            '--headless',
            '--convert-to', 'png',
            abs_pptx_filepath,
            '--outdir', abs_output_folder
        ]
        # Increased timeout to handle large files
        # Snapshot current images to filter only newly generated files
        pre_existing = set(str(p) for p in Path(abs_output_folder).glob('*.png'))
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=120)
        post_existing = set(str(p) for p in Path(abs_output_folder).glob('*.png'))
        newly_created = list(post_existing - pre_existing)
        # If LibreOffice overwrote in-place (no set difference), fall back to all PNGs in the folder
        if not newly_created:
            newly_created = [str(p) for p in Path(abs_output_folder).glob('*.png')]

    except FileNotFoundError:
        # This is a fallback, should not be hit if shutil.which() is accurate.
        print(f"ERROR: The '{libreoffice_exec}' command was not found, even though it was detected.")
        print("Please ensure LibreOffice is installed and that its program directory is in your system's PATH.")
        return []
    except subprocess.TimeoutExpired:
        print("ERROR: Timeout expired while converting slides with LibreOffice. The presentation may be too large or complex.")
        return []
    except Exception as e:
        print(f"Error exporting slides using LibreOffice: {e}")
        # Capture stderr for more detailed error info if available
        if hasattr(e, 'stderr'):
            print(f"LibreOffice stderr: {e.stderr.decode()}")
        return []

    # Find generated files for this conversion only, sort naturally so slide10 comes after slide9
    def get_slide_num(f):
        match = re.search(r'(\d+)', os.path.basename(f))
        return int(match.group(1)) if match else -1

    # Fallback if detection failed
    candidate_paths = newly_created if 'newly_created' in locals() else [str(p) for p in Path(abs_output_folder).glob('*.png')]
    candidate_paths = [p for p in candidate_paths if os.path.exists(p)]
    candidate_paths.sort(key=get_slide_num)

    # If we only got a single PNG, LibreOffice likely exported only the first slide.
    # Fallback: export to PDF and rasterize pages to PNGs using pdftoppm if available.
    if len(candidate_paths) <= 1:
        try:
            pdf_path = convert_pptx_to_pdf(abs_pptx_filepath, abs_output_folder)
            pdftoppm_exec = shutil.which('pdftoppm')
            if pdf_path and pdftoppm_exec:
                prefix = os.path.join(abs_output_folder, 'slide')
                # -png writes slide-1.png, slide-2.png, ...
                subprocess.run([pdftoppm_exec, '-png', pdf_path, prefix], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=120)
                pngs = sorted([str(p) for p in Path(abs_output_folder).glob('slide-*.png')], key=get_slide_num)
                if pngs:
                    return pngs
        except Exception as e:
            print(f"PDF fallback image export failed: {e}")

    return candidate_paths

def convert_pptx_to_pdf(pptx_filepath, output_folder):
    """Converts a .pptx file to PDF using LibreOffice."""
    libreoffice_exec = None
    for name in ['soffice', 'libreoffice']:
        if shutil.which(name):
            libreoffice_exec = name
            break

    if not libreoffice_exec:
        print("ERROR: Could not find LibreOffice for PDF conversion.")
        return None

    abs_pptx_filepath = os.path.abspath(pptx_filepath)
    abs_output_folder = os.path.abspath(output_folder)
    Path(abs_output_folder).mkdir(parents=True, exist_ok=True)

    try:
        cmd = [
            libreoffice_exec,
            '--headless',
            '--convert-to', 'pdf',
            abs_pptx_filepath,
            '--outdir', abs_output_folder
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=120)
        pdf_path = Path(abs_output_folder) / (Path(abs_pptx_filepath).stem + '.pdf')
        if pdf_path.exists():
            return str(pdf_path)
    except Exception as e:
        print(f"Error converting PPTX to PDF: {e}")
    return None

def extract_specific_xml_from_pptx(pptx_filepath, xml_filename):
    """
    Extracts the content of a single specified XML file from a .pptx file.
    
    Args:
        pptx_filepath (str): Path to the .pptx file.
        xml_filename (str): The internal path to the XML file (e.g., 'ppt/slides/slide1.xml').
        
    Returns:
        str: The content of the XML file as a string, or None if not found.
    """
    try:
        with zipfile.ZipFile(pptx_filepath, 'r') as pptx_zip:
            # Normalize filename for matching
            xml_filename_normalized = xml_filename.replace("\\", "/")
            if xml_filename_normalized in pptx_zip.namelist():
                with pptx_zip.open(xml_filename_normalized) as xml_file:
                    return xml_file.read().decode('utf-8')
            return None
    except Exception as e:
        print(f"Error extracting specific XML '{xml_filename}' from {pptx_filepath}: {e}")
        return None

def image_to_base64(image_path):
    """Converts an image file to a base64 encoded string."""
    try:
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    except Exception as e:
        print(f"Error converting image {image_path} to base64: {e}")
        return None


def validate_xml(xml_text):
    """Return True if xml_text is well-formed."""
    try:
        from lxml import etree
        etree.fromstring(xml_text.encode("utf-8"))
        return True
    except Exception as e:
        print(f"Invalid XML detected: {e}")
        return False

def attempt_repair_xml(xml_text: str) -> str | None:
    """
    Attempt to repair malformed XML using lxml's recovery parser.
    Returns a repaired XML string if successful, otherwise None.
    """
    try:
        from lxml import etree
        parser = etree.XMLParser(recover=True)
        root = etree.fromstring(xml_text.encode("utf-8"), parser=parser)
        if root is None:
            return None
        # Serialize back to UTF-8 string
        fixed = etree.tostring(root, encoding="utf-8", xml_declaration=True).decode("utf-8")
        return fixed
    except Exception as e:
        print(f"XML repair failed: {e}")
        return None


def convert_pptx_to_base64_images(pptx_filepath):
    """Converts each slide of a .pptx file to base64 encoded PNG images."""
    temp_dir = tempfile.mkdtemp()
    images = export_slides_to_images(pptx_filepath, temp_dir)
    base64_images = []
    for img_path in images:
        try:
            with open(img_path, 'rb') as f:
                img_str = base64.b64encode(f.read()).decode('utf-8')
                base64_images.append(f"data:image/png;base64,{img_str}")
        except Exception as e:
            print(f"Error encoding image {img_path} to base64: {e}")
    shutil.rmtree(temp_dir, ignore_errors=True)
    return base64_images
